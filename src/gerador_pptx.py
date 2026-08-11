from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

from .leitor_txt import ler_txt, dividir_em_telas

def _clonar_slide(prs_destino, slide_origem):
    novo = prs_destino.slides.add_slide(prs_destino.slide_layouts[6])

    # Remove placeholders padrão do slide vazio
    for shape in list(novo.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    for shape in slide_origem.shapes:
        novo.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")

    return novo

def _formas_texto(slide):
    formas = []
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.has_text_frame:
            formas.append(shape)
    return formas

def _definir_texto(slide, texto):
    formas = _formas_texto(slide)
    if not formas:
        raise ValueError("O slide do modelo não possui caixa de texto editável.")

    # Prioriza a maior caixa de texto
    shape = max(formas, key=lambda s: int(s.width) * int(s.height))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    partes = texto.split("\n")
    p = tf.paragraphs[0]
    p.text = partes[0] if partes else ""
    for linha in partes[1:]:
        novo_p = tf.add_paragraph()
        novo_p.text = linha

def gerar_pptx(caminho_txt, pasta_saida, modelo_pptx):
    caminho_txt = Path(caminho_txt)
    pasta_saida = Path(pasta_saida)
    modelo_pptx = Path(modelo_pptx)

    if not modelo_pptx.is_file():
        raise FileNotFoundError("O modelo PowerPoint selecionado não foi encontrado.")

    titulo, linhas = ler_txt(caminho_txt)
    telas = dividir_em_telas(linhas)

    modelo = Presentation(modelo_pptx)
    if len(modelo.slides) < 2:
        raise ValueError(
            "O modelo PowerPoint precisa ter pelo menos 2 slides: "
            "slide 1 para título e slide 2 para letra."
        )

    # Cria uma apresentação nova mantendo dimensões do modelo
    saida = Presentation()
    saida.slide_width = modelo.slide_width
    saida.slide_height = modelo.slide_height

    # Remove slide inicial, se existir
    while len(saida.slides):
        rid = saida.slides._sldIdLst[0].rId
        saida.part.drop_rel(rid)
        del saida.slides._sldIdLst[0]

    capa = _clonar_slide(saida, modelo.slides[0])
    _definir_texto(capa, titulo)

    for tela in telas:
        slide = _clonar_slide(saida, modelo.slides[1])
        _definir_texto(slide, "\n".join(tela))

    pasta_saida.mkdir(parents=True, exist_ok=True)
    destino = pasta_saida / f"{caminho_txt.stem}.pptx"
    saida.save(destino)
    return destino
